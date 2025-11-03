"""
Generate PowerPoint presentation for classification UQ results
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_title_slide(prs, title, subtitle):
    """Create title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide

def create_content_slide(prs, title):
    """Create content slide with title and content placeholder"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    return slide

def add_bullet_points(textbox, points):
    """Add bullet points to textbox"""
    tf = textbox.text_frame
    tf.clear()
    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)

def create_table_slide(prs, title, data, headers):
    """Create slide with table"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    
    # Title
    title_box = slide.shapes.title
    title_box.text = title
    
    # Table
    rows = len(data) + 1  # +1 for header
    cols = len(headers)
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(0.5 * rows)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Header row
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Data rows
    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, value in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(12)
    
    return slide

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    create_title_slide(
        prs,
        "Uncertainty Quantification in Medical Image Classification",
        "Extension to Multi-Class Classification Tasks"
    )
    
    # Slide 2: Objectives
    slide = create_content_slide(prs, "Project Update: Objectives")
    content = slide.placeholders[1]
    add_bullet_points(content, [
        "Extend UQ methods from segmentation to multi-class classification",
        "Evaluate performance on Chest X-Ray Pneumonia dataset",
        "Compare Baseline, MC Dropout, Deep Ensemble, and SWAG methods",
        "Assess calibration (ECE) and probabilistic quality (Brier score)",
        "Analyze trade-offs between accuracy and uncertainty quantification"
    ])
    
    # Slide 3: Dataset
    slide = create_content_slide(prs, "Dataset: Chest X-Ray Pneumonia")
    content = slide.placeholders[1]
    add_bullet_points(content, [
        "Binary classification: Normal vs. Pneumonia",
        "Training: 4,172 images (augmented via calibration split)",
        "Calibration: 1,044 images (20% of training for CRC)",
        "Test: 624 images",
        "Image size: 224×224 (resized from original chest X-rays)",
        "Architecture: ResNet-18 (11.2M parameters)",
        "Training: 5 epochs with Adam optimizer, CPU"
    ])
    
    # Slide 4: Evaluation Results
    results_data = [
        ["Baseline", "85.26%", "0.0617", "0.2445", "—"],
        ["MC Dropout (T=20)", "86.38%", "0.0538", "0.2074", "0.0017"],
        ["Deep Ensemble (M=5)", "88.78%", "0.0340", "0.1727", "0.0371"],
        ["SWAG (T=30)", "82.85%", "0.0996", "0.2745", "—"]
    ]
    headers = ["Method", "Accuracy", "ECE ↓", "Brier Score ↓", "Mean Uncertainty"]
    create_table_slide(prs, "Classification Results Summary", results_data, headers)
    
    # Slide 5: SWAG Overview
    slide = create_content_slide(prs, "SWAG: Stochastic Weight Averaging-Gaussian")
    content = slide.placeholders[1]
    add_bullet_points(content, [
        "Bayesian approximation via weight-space sampling",
        "Collects model snapshots after epoch 3 (post-annealing)",
        "Maintains running mean and covariance of weights",
        "Samples T=30 models from posterior for prediction",
        "Trade-off: Lower accuracy (82.85%) but richer uncertainty",
        "Higher ECE (0.0996) suggests miscalibration in this setup",
        "Potential: Better uncertainty with longer training/tuning"
    ])
    
    # Slide 6: Key Findings
    slide = create_content_slide(prs, "Key Findings")
    content = slide.placeholders[1]
    add_bullet_points(content, [
        "Deep Ensemble achieves best accuracy (88.78%) and calibration (ECE=0.034)",
        "MC Dropout improves over baseline with minimal overhead",
        "SWAG underperforms in short training regime (5 epochs)",
        "Calibration metrics (ECE, Brier) correlate with ensemble diversity",
        "Classification tasks show clearer UQ benefits than segmentation"
    ])
    
    # Slide 7: Comparison: Classification vs. Segmentation
    slide = create_content_slide(prs, "Classification vs. Segmentation Performance")
    content = slide.placeholders[1]
    add_bullet_points(content, [
        "Classification: Higher baseline accuracy (85% vs. ~70% Dice)",
        "Ensemble improvement: +3.5% accuracy (classification) vs. +5% Dice (segmentation)",
        "Calibration: Easier to measure in classification (single prob per image)",
        "Uncertainty quality: More interpretable in classification (class probs)",
        "Runtime: Classification inference ~10× faster per sample",
        "Limitation: Classification loses spatial uncertainty (segmentation strength)"
    ])
    
    # Slide 8: Limitations & Future Work
    slide = create_content_slide(prs, "Limitations & Future Work")
    content = slide.placeholders[1]
    add_bullet_points(content, [
        "Short training (5 epochs) limits SWAG performance",
        "Conformal Risk Control not yet integrated (time constraint)",
        "CPU-only runs → slower iteration, smaller batch sizes",
        "Future: Extend to multi-class (>2 classes), test on OCT/Brain Tumor",
        "Future: Integrate CRC for FNR control in medical diagnosis",
        "Future: GPU training for deeper models and hyperparameter tuning"
    ])
    
    # Save presentation
    output_path = 'presentation/classification_results.pptx'
    prs.save(output_path)
    print(f'Presentation saved to {output_path}')

if __name__ == '__main__':
    main()
