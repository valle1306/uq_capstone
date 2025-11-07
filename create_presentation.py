"""
Create a comprehensive PPTX presentation for UQ Capstone Project
Medical Image Classification with Uncertainty Quantification Methods
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import json
from pathlib import Path

# Initialize presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
COLOR_PRIMARY = RGBColor(0, 102, 204)      # Blue
COLOR_ACCENT = RGBColor(220, 20, 60)       # Crimson
COLOR_SUCCESS = RGBColor(34, 139, 34)      # Forest Green
COLOR_WARNING = RGBColor(255, 140, 0)      # Dark Orange
COLOR_TEXT = RGBColor(40, 40, 40)          # Dark Gray
COLOR_LIGHT = RGBColor(245, 245, 245)      # Light Gray

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRIMARY
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items=None, add_image=None):
    """Add content slide with bullet points or image"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Add title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLOR_PRIMARY
    title_shape.line.color.rgb = COLOR_PRIMARY
    
    # Add title text
    title_frame = title_shape.text_frame
    title_frame.word_wrap = False
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.space_before = Pt(10)
    p.space_after = Pt(10)
    
    # Add content
    if add_image:
        # Add image
        try:
            slide.shapes.add_picture(add_image, Inches(0.5), Inches(1.2), width=Inches(9))
        except:
            pass
    elif content_items:
        # Add bullet points
        text_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(6))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        for i, item in enumerate(content_items):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            if isinstance(item, tuple):
                text, level = item
            else:
                text = item
                level = 0
            
            p.text = text
            p.level = level
            p.font.size = Pt(22 - level * 3)
            p.font.color.rgb = COLOR_TEXT
            p.space_before = Pt(6)
            p.space_after = Pt(6)
    
    return slide

def add_table_slide(prs, title, table_data):
    """Add slide with table"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Add title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLOR_PRIMARY
    title_shape.line.color.rgb = COLOR_PRIMARY
    
    title_frame = title_shape.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.space_before = Pt(10)
    p.space_after = Pt(10)
    
    # Add table
    rows, cols = len(table_data), len(table_data[0])
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(9)
    height = Inches(5.5)
    
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Fill table
    for i, row in enumerate(table_data):
        for j, cell_text in enumerate(row):
            cell = table_shape.cell(i, j)
            cell.text = str(cell_text)
            
            # Format header row
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_PRIMARY
                p = cell.text_frame.paragraphs[0]
                p.font.bold = True
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(255, 255, 255)
            else:
                if i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLOR_LIGHT
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(12)
                p.font.color.rgb = COLOR_TEXT
            
            p.alignment = PP_ALIGN.CENTER
    
    return slide

# ============= SLIDE 1: Title Slide =============
add_title_slide(
    prs,
    "Uncertainty Quantification for Medical Image Classification",
    "Capstone Project - Rutgers University\nNovember 2025"
)

# ============= SLIDE 2: Problem Statement =============
add_content_slide(
    prs,
    "Problem Statement",
    [
        ("Challenge: Medical AI demands reliable uncertainty estimates", 0),
        ("Why? Model confidence doesn't equal correctness", 0),
        ("", 0),
        ("Real-world scenarios:", 0),
        ("A model says 92% confident in chest X-ray diagnosis", 1),
        ("But what's the actual error rate at that confidence?", 1),
        ("How do we know when to trust model predictions?", 1),
        ("", 0),
        ("Solution: Compare Uncertainty Quantification (UQ) methods", 0),
        ("Evaluate 4 different UQ approaches on same dataset", 1),
        ("Measure accuracy, calibration, and uncertainty quality", 1),
        ("Identify best method for medical imaging applications", 1),
    ]
)

# ============= SLIDE 3: Dataset Overview =============
add_content_slide(
    prs,
    "Dataset: Chest X-Ray Binary Classification",
    [
        ("Source: Kaggle Chest X-Ray Images (Pneumonia)", 0),
        ("", 0),
        ("Split:", 0),
        ("Training: 4,172 images", 1),
        ("Validation (Calibration): 1,044 images", 1),
        ("Test: 624 images", 1),
        ("", 0),
        ("Task: Binary classification - Normal vs Pneumonia", 0),
        ("", 0),
        ("Model Architecture: ResNet-18 pretrained on ImageNet", 0),
        ("Baseline accuracy: 91.67%", 1),
    ]
)

# ============= SLIDE 4: UQ Methods Overview =============
add_content_slide(
    prs,
    "Uncertainty Quantification Methods",
    [
        ("1. Baseline (No Uncertainty)", 0),
        ("   Standard ResNet-18 with cross-entropy loss", 1),
        ("   Provides only point estimate of class", 1),
        ("", 0),
        ("2. MC Dropout (Bayesian Approximation)", 0),
        ("   Dropout during inference (T=15 samples)", 1),
        ("   Uncertainty from prediction variance", 1),
        ("", 0),
        ("3. Deep Ensemble (Ensemble Method)", 0),
        ("   5 independently trained models", 1),
        ("   Uncertainty from model disagreement", 1),
        ("", 0),
        ("4. SWAG (Stochastic Weight Averaging Gaussian)", 0),
        ("   Bayesian posterior approximation", 1),
        ("   30 samples from weight distribution", 1),
    ]
)

# ============= SLIDE 5: Results - Accuracy Comparison =============
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

# Title bar
title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
title_shape.fill.solid()
title_shape.fill.fore_color.rgb = COLOR_PRIMARY
title_frame = title_shape.text_frame
p = title_frame.paragraphs[0]
p.text = "Results: Accuracy Comparison"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.space_before = Pt(10)

# Add image
try:
    slide.shapes.add_picture(
        "runs/classification/metrics/comprehensive_metrics_visualization.png",
        Inches(0.5), Inches(1.0), width=Inches(9)
    )
except:
    text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    text_frame = text_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = "Visualizations available in: runs/classification/metrics/"
    p.font.size = Pt(20)

# ============= SLIDE 6: Results - Summary Table =============
add_table_slide(
    prs,
    "Results Summary",
    [
        ["Method", "Accuracy", "ECE", "Brier", "FNR"],
        ["Baseline", "91.67%", "0.0498", "0.0704", "0.0833"],
        ["MC Dropout", "85.26%", "0.1172", "0.1246", "0.1474"],
        ["Deep Ensemble", "91.67%", "0.0271", "0.0630", "0.0833"],
        ["SWAG", "83.17%", "0.1519", "0.1528", "0.1683"],
    ]
)

# ============= SLIDE 7: Analysis - MC Dropout Success Story =============
add_content_slide(
    prs,
    "MC Dropout: Debugging Success",
    [
        ("Initial Issue: MC Dropout evaluated at 66% accuracy", 0),
        ("   Expected: ~90% (competitive with baseline)", 1),
        ("   Problem: Severe evaluation bug", 1),
        ("", 0),
        ("Root Cause Found:", 0),
        ("   Dropout was enabled ONCE before all samples", 1),
        ("   Accumulated across T=15 forward passes", 1),
        ("   Corrupted predictions & destroyed performance", 1),
        ("", 0),
        ("Solution: Toggle dropout on/off between each sample", 0),
        ("   model.enable_dropout() → forward() → model.eval()", 1),
        ("", 0),
        ("Result: 66% → 85.26% (FIXED!)", 0),
        ("   Demonstrates importance of careful evaluation protocol", 1),
    ]
)

# ============= SLIDE 8: Analysis - SWAG Underperformance =============
add_content_slide(
    prs,
    "SWAG: Why Only 83.17%?",
    [
        ("Expected: ~91% (competitive with Ensemble)", 0),
        ("Actual: 83.17% (underperformer)", 0),
        ("", 0),
        ("Root Cause: Validation Overfitting", 0),
        ("   Training metrics: 99.62% validation accuracy", 1),
        ("   But test accuracy: 85.58% (training phase)", 1),
        ("   Posterior sampling: 83.17% (evaluation)", 1),
        ("", 0),
        ("Why overfitting happened:", 0),
        ("   Retrained from baseline with same hyperparameters", 1),
        ("   Snapshot collection epochs 30-50 caught overfit region", 1),
        ("   Weight distribution skewed toward overfit solution", 1),
        ("", 0),
        ("Key insight:", 0),
        ("   Bayesian posterior sampling naturally includes uncertainty", 1),
        ("   When posterior is concentrated on overfit solution → poor OOD", 1),
    ]
)

# ============= SLIDE 9: Ensemble Superiority =============
add_content_slide(
    prs,
    "Winner: Deep Ensemble",
    [
        ("Accuracy: 91.67% (tied with Baseline on test set)", 0),
        ("Calibration: ECE = 0.0271 (BEST)", 0),
        ("Brier Score: 0.0630 (excellent)", 0),
        ("", 0),
        ("Why Ensemble Dominates:", 0),
        ("", 0),
        ("1. Model Diversity", 0),
        ("   5 independent initializations capture different perspectives", 1),
        ("", 0),
        ("2. Robust Predictions", 0),
        ("   Disagreement = epistemic uncertainty", 1),
        ("   Majority voting = better calibration", 1),
        ("", 0),
        ("3. Simple & Interpretable", 0),
        ("   No complex weight distributions", 1),
        ("   Easy to implement in production", 1),
        ("", 0),
        ("Trade-off: 5x computational cost at inference", 0),
    ]
)

# ============= SLIDE 10: Method Comparison =============
add_content_slide(
    prs,
    "Method Comparison: Pros & Cons",
    [
        ("MC Dropout (85.26%)", 0),
        ("   Pros: Fast inference, minimal memory, provides uncertainty", 1),
        ("   Cons: Lower accuracy, weaker calibration", 1),
        ("", 0),
        ("Deep Ensemble (91.67%)", 0),
        ("   Pros: Best accuracy, best calibration, interpretable", 1),
        ("   Cons: 5x inference cost, requires 5 model copies", 1),
        ("", 0),
        ("SWAG (83.17%)", 0),
        ("   Pros: Theoretically sound Bayesian method", 1),
        ("   Cons: Sensitive to hyperparameters, overfitting issues", 1),
        ("   Requires careful regularization tuning", 1),
        ("", 0),
        ("Recommendation for Medical Imaging:", 0),
        ("   Use Ensemble if computational budget allows", 1),
        ("   Use MC Dropout if speed critical", 1),
    ]
)

# ============= SLIDE 11: Conformal Risk Control =============
add_content_slide(
    prs,
    "Conformal Risk Control: Post-hoc Calibration",
    [
        ("What: Distribution-free confidence guarantees", 0),
        ("", 0),
        ("Goal: Generate prediction sets with specified coverage", 0),
        ("   Example: 95% confidence set contains true label", 1),
        ("   Guaranteed even under distribution shift", 1),
        ("", 0),
        ("Methods Implemented:", 0),
        ("", 0),
        ("1. FNR Control (α=0.05, 0.10)", 0),
        ("   Controls False Negative Rate at calibration level", 1),
        ("   Critical for medical applications (avoid missed positives)", 1),
        ("", 0),
        ("2. Set Size Control", 0),
        ("   Minimizes prediction set size while maintaining coverage", 1),
        ("", 0),
        ("3. Composite Loss", 0),
        ("   Balances coverage and set size", 1),
        ("", 0),
        ("Why Important for Medical AI:", 0),
        ("   Uncertainty quantified with formal guarantees", 1),
        ("   Robust to domain shift between hospitals", 1),
    ]
)

# ============= SLIDE 12: Segmentation Comparison (Context) =============
add_content_slide(
    prs,
    "Context: Segmentation vs Classification",
    [
        ("Previous Work: UQ for Chest X-Ray Segmentation", 0),
        ("", 0),
        ("Classification Results (This Project):", 0),
        ("   Best Method: Deep Ensemble (91.67%)", 1),
        ("   Runner-up: Baseline (91.67%)", 1),
        ("   MC Dropout: 85.26%", 1),
        ("   SWAG: 83.17%", 1),
        ("", 0),
        ("Cross-Task Pattern:", 0),
        ("   Ensemble consistently strong across both tasks", 1),
        ("   MC Dropout shows potential but needs tuning", 1),
        ("   SWAG struggles with regularization", 1),
        ("   → Ensemble is go-to for medical imaging", 1),
        ("", 0),
        ("Lesson Learned:", 0),
        ("   Task-agnostic UQ methods benefit from diversity", 1),
        ("   Bayesian methods need careful hyperparameter tuning", 1),
    ]
)

# ============= SLIDE 13: Lessons & Takeaways =============
add_content_slide(
    prs,
    "Key Lessons",
    [
        ("1. Implementation Details Matter", 0),
        ("   MC Dropout dropout toggle was critical (66% → 85%)", 1),
        ("   Careful evaluation protocol essential", 1),
        ("", 0),
        ("2. Overfitting Affects UQ Methods Differently", 0),
        ("   Bayesian methods concentrate on overfit solution", 1),
        ("   Ensemble diversity mitigates this problem", 1),
        ("", 0),
        ("3. Calibration ≠ Accuracy", 0),
        ("   SWAG lower accuracy BUT could calibrate well", 1),
        ("   ECE (Expected Calibration Error) is different metric", 1),
        ("   Both matter for reliable medical AI", 1),
        ("", 0),
        ("4. Conformal Methods Provide Formal Guarantees", 0),
        ("   Post-hoc calibration without model retraining", 1),
        ("   Distribution-free guarantees valuable for deployment", 1),
    ]
)

# ============= SLIDE 14: Next Steps & Future Work =============
add_content_slide(
    prs,
    "Next Steps & Future Directions",
    [
        ("Short-term (Project Completion):", 0),
        ("   Finalize conformal prediction implementation", 1),
        ("   Generate deployment-ready baseline comparisons", 1),
        ("", 0),
        ("Medium-term (Course Extension):", 0),
        ("   Improve SWAG with better regularization", 1),
        ("   Explore adaptive SWAG snapshot collection", 1),
        ("   Test on out-of-distribution data", 1),
        ("", 0),
        ("Research Directions:", 0),
        ("   Combine Ensemble + Conformal for certified safety", 1),
        ("   Investigate task-aware UQ methods", 1),
        ("   Study UQ under realistic distribution shifts", 1),
        ("   Develop explainability for uncertain predictions", 1),
    ]
)

# ============= SLIDE 15: Research Collaboration Proposal =============
add_content_slide(
    prs,
    "Research Collaboration with Dr. Gemma Moran",
    [
        ("Background:", 0),
        ("   Biomathematics + Data Science: You", 1),
        ("   Statistical Learning Theory: Dr. Moran", 1),
        ("", 0),
        ("Proposed Collaboration Themes:", 0),
        ("", 0),
        ("1. Adaptive UQ for Medical Imaging", 0),
        ("   Novel method: Task-aware uncertainty weighting", 1),
        ("   Combine Bayesian + Ensemble + Conformal", 1),
        ("   Original contribution (not yet published)", 1),
        ("", 0),
        ("2. UQ under Distribution Shift", 0),
        ("   Problem: Hospital imaging protocols vary", 1),
        ("   Solution: Robust UQ guarantees across domains", 1),
        ("   Application-driven theory paper", 1),
    ]
)

# ============= SLIDE 16: Research Collaboration Details =============
add_content_slide(
    prs,
    "Research Idea: Adaptive Bayesian Ensembles",
    [
        ("Core Idea:", 0),
        ("   Not all uncertainty sources equally important", 1),
        ("   Medical imaging: Epistemic >> Aleatoric", 1),
        ("   Standard methods treat both the same", 1),
        ("", 0),
        ("Method:", 0),
        ("   Learn task-specific weights for UQ components", 1),
        ("   Ensemble variance (epistemic) vs MC Dropout (aleatoric)", 1),
        ("   Optimize for clinical relevance (calibration + decision-making)", 1),
        ("", 0),
        ("Why Novel:", 0),
        ("   Combines multiple UQ paradigms adaptively", 1),
        ("   No published work on this specific combination", 1),
        ("   Bridges biomathematical modeling + practical ML", 1),
        ("", 0),
        ("Your Role:", 0),
        ("   Mathematical framework + empirical validation", 1),
        ("   Dr. Moran: Statistical theory + formal guarantees", 1),
    ]
)

# ============= SLIDE 17: Alternative Collaboration Ideas =============
add_content_slide(
    prs,
    "Alternative Research Directions",
    [
        ("Option 2: Conformal Regression for Risk Prediction", 0),
        ("   Extend conformal methods to severity scoring", 1),
        ("   Provide prediction intervals with coverage guarantees", 1),
        ("   Application: Disease progression forecasting", 1),
        ("", 0),
        ("Option 3: UQ Quality Metrics Framework", 0),
        ("   Formal definition of 'good' uncertainty", 1),
        ("   Develop diagnostics for UQ methods", 1),
        ("   Publication: \"Uncertainty Quality in Medical AI\"", 1),
        ("", 0),
        ("Option 4: Interpretable UQ for Clinicians", 0),
        ("   Make uncertainty explanations actionable", 1),
        ("   What features cause model uncertainty?", 1),
        ("   Collaboration with hospital domain experts", 1),
        ("", 0),
        ("Common Theme:", 0),
        ("   Practical, not overly theoretical", 1),
        ("   Medical imaging focus", 1),
        ("   Publication-ready original contribution", 1),
    ]
)

# ============= SLIDE 18: Timeline & Collaboration Structure =============
add_content_slide(
    prs,
    "Proposed Timeline & Collaboration",
    [
        ("Phase 1 (Now - December):", 0),
        ("   Finalize current capstone project", 1),
        ("   Literature review on adaptive UQ methods", 1),
        ("   Initial experiments on your dataset", 1),
        ("", 0),
        ("Phase 2 (January - March):", 0),
        ("   Develop theoretical framework with Dr. Moran", 1),
        ("   Implement adaptive method", 1),
        ("   Comprehensive comparison on multiple datasets", 1),
        ("", 0),
        ("Phase 3 (April - May):", 0),
        ("   Write manuscript", 1),
        ("   Submit to medical imaging conference (MICCAI, ISBI, etc.)", 1),
        ("", 0),
        ("Deliverables:", 0),
        ("   1-2 first-author conference papers", 1),
        ("   Open-source implementation", 1),
        ("   Medical imaging community contribution", 1),
    ]
)

# ============= SLIDE 19: Your Unique Position =============
add_content_slide(
    prs,
    "Your Unique Advantage",
    [
        ("Biomathematics Background:", 0),
        ("   Comfortable with mathematical rigor", 1),
        ("   Understanding of biological constraints", 1),
        ("", 0),
        ("Data Science Skills:", 0),
        ("   Practical ML implementation expertise", 1),
        ("   Familiar with medical datasets", 1),
        ("   This capstone project is proof of capability", 1),
        ("", 0),
        ("Why This Collaboration Works:", 0),
        ("   You bring applied perspective to statistical theory", 1),
        ("   Dr. Moran brings rigor to practical methods", 1),
        ("   Perfect match for modern ML research", 1),
        ("", 0),
        ("Current Moment:", 0),
        ("   UQ in medical AI is hot topic", 1),
        ("   Your work shows genuine research potential", 1),
        ("   Now is the time to expand into publications", 1),
    ]
)

# ============= SLIDE 20: Conclusion =============
add_title_slide(
    prs,
    "Questions?",
    "Thank you for your attention\nContact: [Your Email]\nGitHub: https://github.com/valle1306/uq_capstone"
)

# Save presentation
output_path = "presentation/UQ_Capstone_Presentation.pptx"
Path("presentation").mkdir(exist_ok=True)
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
